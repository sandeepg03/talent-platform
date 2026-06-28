"""
generate_submission_materials.py

Generates both the presentation files required for the hackathon submission:
1. AI_Talent_Intelligence_Platform.pptx (using python-pptx)
2. AI_Talent_Intelligence_Platform.pdf (using reportlab)

Contains slide structures for:
- Problem Statement, Proposed Solution, System Architecture, Tech Stack, Dataset Overview,
  Workflow/Pipeline, Key Features, AI/ML Approach, Scoring & Ranking Methodology,
  Scalability, Challenges Faced, Results, Future Improvements, Demo, Conclusion.
"""

from __future__ import annotations

import os
import sys

# --- PPTX GENERATION ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- PDF GENERATION ---
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

# Cohesive Color Palette
DARK_NAVY = RGBColor(11, 25, 44)
TEXT_WHITE = RGBColor(255, 255, 255)
PRIMARY_BLUE = RGBColor(27, 38, 59)
SECONDARY_BLUE = RGBColor(65, 90, 119)
LIGHT_GRAY = RGBColor(244, 245, 247)
ACCENT_CYAN = RGBColor(0, 180, 216)
MUTED_DARK = RGBColor(51, 65, 85)

PDF_DARK_NAVY = colors.HexColor("#0B192C")
PDF_PRIMARY_BLUE = colors.HexColor("#1B263B")
PDF_SECONDARY_BLUE = colors.HexColor("#415A77")
PDF_LIGHT_GRAY = colors.HexColor("#F4F5F7")
PDF_ACCENT_CYAN = colors.HexColor("#00B4D8")
PDF_MUTED_DARK = colors.HexColor("#334155")

SLIDE_DATA = [
    {
        "title": "AI Talent Intelligence Platform",
        "subtitle": "Production-Grade Semantic Search, Reranking & Hybrid Scoring Pipeline for Large-Scale Candidate Ranking (100K+ Dataset)",
        "points": [
            "Hackathon Track: Artificial Intelligence & Data Science",
            "Author Profile: sandeepg03",
            "GitHub Repository: https://github.com/sandeepg03/talent-platform",
            "Core Pipeline: Bi-Encoder (FAISS) -> Cross-Encoder Reranker -> Feature Engineer -> Hybrid Scorer"
        ]
    },
    {
        "title": "Problem Statement",
        "subtitle": "Bottlenecks in Modern Resume Parsing and Ranking Pipelines",
        "points": [
            "Scale Constraints: Processing 100K+ resumes in real-time is computationally heavy; naive O(N2) similarity checks time out.",
            "Weak Keyword Matching: Standard parser matching misses semantic context and synonyms (e.g., matching 'ML' with 'Machine Learning').",
            "Lack of Transparency: Most modern AI ranking systems function as 'black boxes' without clear, explainable reasoning for rankings.",
            "Adversarial Resumes (Honeypots): Systems are susceptible to keyword stuffing and fake profiles designed to cheat metrics."
        ]
    },
    {
        "title": "Proposed Solution",
        "subtitle": "Two-Stage Cascade Architecture with Transparent Hybrid Scoring",
        "points": [
            "Cascade Search Pipeline: High-recall bi-encoder fast retrieval followed by high-precision cross-encoder reranking.",
            "Explainable AI: Every rank has a deterministic, recruiter-friendly generated explanation detailing skills, experience, and signal alignment.",
            "Honeypot Shielding: Integrated heuristics filter out bad-faith synthetic profiles at the feature-engineering layer.",
            "Flexible Hybrid Scorer: Balanced combination of semantic similarity, deep cross-encoder matching, and recruiter-specific heuristics."
        ]
    },
    {
        "title": "System Architecture",
        "subtitle": "Modular End-to-End Cascade Architecture",
        "points": [
            "JD Parser: Extracts structured criteria (must-have/nice-to-have skills, years of experience, target domains).",
            "Embedding Stage: Encodes JD and candidate text into 384-dim dense vectors using BAAI/bge-small-en-v1.5.",
            "FAISS Vector Store: Performs sub-5ms IndexFlatIP inner-product search to extract top-500 candidate pool.",
            "Cross-Encoder Reranking: Computes deep query-resume attention on top-500 using ms-marco-MiniLM-L-6-v2, keeping top-200.",
            "Feature Engineering: Extracts years of experience, education level, certificate matches, open-source signals, and flags honeypots.",
            "Hybrid Scoring & Reasoning: Generates final weighted score (0-100) and writes detailed recruiter-centric explanation."
        ]
    },
    {
        "title": "Technology Stack",
        "subtitle": "Modern, Production-Grade open-source AI & Data Toolkit",
        "points": [
            "Embedding & Retrieval: Sentence-Transformers (BGE-Small-en-v1.5) & Facebook AI Similarity Search (FAISS).",
            "Deep Learning Reranking: PyTorch & Cross-Encoder Rerankers (ms-marco-MiniLM-L-6-v2).",
            "Schema Validation & Models: Pydantic v2 for strict type safety and structured schemas.",
            "Web APIs & Servers: FastAPI backend for ultra-low latency server requests & Uvicorn ASGI server.",
            "Frontend App: Streamlit dashboard with interactive Plotly visualization charts.",
            "Testing & CI/CD: Pytest framework (362 tests) & Docker + Compose + GitHub Actions CI workflows."
        ]
    },
    {
        "title": "Dataset Overview",
        "subtitle": "100,000+ Structured Candidates and Target Job Descriptions",
        "points": [
            "Scale: 100,000+ candidate profiles in JSONL format (~464MB raw text).",
            "Structured Profiles: Fully typed schema including experience, skills, career history, education, certifications, and GitHub signals.",
            "Adversarial Profiles (Honeypots): Injected synthetic profiles designed to exploit simple keyword engines.",
            "Target JD: Word Document (.docx) detailing structured target roles, must-have skills, and nice-to-have capabilities."
        ]
    },
    {
        "title": "Workflow / Pipeline",
        "subtitle": "Staggered Offline Precomputation vs Low-Latency Online Ranking",
        "points": [
            "Offline Phase (precompute.py):",
            "  - Stream-parses 100k candidates from JSONL (9.5s). Truncates profiles to 96 words (matches model token limits).",
            "  - Staggered multiprocessing pool (3 workers, 2 threads each) encodes 100k texts in ~15 minutes on CPU.",
            "  - Persists embeddings and builds FAISS index flat-IP file to local disk artifacts.",
            "Online Phase (rank.py):",
            "  - Parses job description -> retrieves top-500 via FAISS index lookup (<5ms).",
            "  - Reranks top-500 to top-200 via cross-encoder -> Scores top-100 -> Writes explanation -> Outputs submission.csv."
        ]
    },
    {
        "title": "Key Features",
        "subtitle": "Advanced Capabilities Designed for Recruiter Workflow",
        "points": [
            "Honeypot Resilience: Automatically scores but filters out adversarial resumes using 4 strict heuristic rules.",
            "Recruiter-Signal Upweighting: Skills with Expert/Advanced proficiency are repeated during text building to gain higher embedding weight.",
            "Deterministic Explainability: Recruiter-friendly reasoning details exact skill matches, experience gap, and signal scores.",
            "Low-Latency API: Low-latency endpoint (/rank) computes, reranks, scores, and filters 100k candidates in seconds.",
            "Interactive Streamlit UI: Upload JDs, browse candidates, inspect score breakdowns, and visualize skill alignments."
        ]
    },
    {
        "title": "AI/ML Approach",
        "subtitle": "State-of-the-Art Dense Retrieval and Cross-Attention Reranking",
        "points": [
            "Bi-Encoder (bge-small-en-v1.5): Top-performing lightweight model on MTEB. Maps queries and candidates to unit-normalized space.",
            "FAISS Indexing: Fast inner product search (L2-norm vectors mean dot-product is equivalent to cosine similarity).",
            "Cross-Encoder (ms-marco-MiniLM-L-6-v2): Computes self-attention across combined query-document tokens.",
            "Context Truncation: Restricting candidate texts to 96 words removes tail noise, matches model limits, and speeds up CPU speed by 4x."
        ]
    },
    {
        "title": "Scoring & Ranking Methodology",
        "subtitle": "Multi-Dimensional Transparent Weighted Scoring System",
        "points": [
            "Weighted Score Formula (Normalized to [0, 100]):",
            "  - 40% Semantic Similarity (Bi-Encoder FAISS cosine score)",
            "  - 30% Cross-Encoder Reranker Score (Sigmoid-normalized logit)",
            "  - 10% Years of Experience Score (Matched against JD target)",
            "  - 10% GitHub Open Source Signal (Activity score)",
            "  - 5% Education Alignment Score (Degree levels matched to JD)",
            "  - 5% Certification Match Score (Exact name matching)",
            "Tie-Breaking: Scores rounded to 4 decimals. Ties resolved by candidate_id ascending (alphabetical order)."
        ]
    },
    {
        "title": "Scalability & Performance",
        "subtitle": "CPU-Optimized Scaling and Low Footprint Memory Design",
        "points": [
            "Staggered Multiprocessing: Avoids OpenMP deadlocks on Windows CPU and eliminates cache/memory bandwidth thrashing.",
            "Streaming Parser: Iterator streams candidates one-by-one, keeping RAM usage under 1.5GB during 100k candidate parsing.",
            "Vector Space Compression: 100k vectors stored in float32 consumes only 150MB of RAM, making it perfectly runnable on local laptops.",
            "Cascaded Filtering: Restricting heavy neural computations to subset layers (100k -> 500 -> 200 -> 100) ensures low latency."
        ]
    },
    {
        "title": "Challenges Faced & Mitigations",
        "subtitle": "Engineering Solutions to Real-World Hardware and System Constraints",
        "points": [
            "1. CPU Encoding Bottleneck: Running 100k long texts took 3+ hours on CPU. Resolved by pre-truncating to 96 words (4x speedup).",
            "2. OpenMP Deadlocks: PyTorch multiprocessing hung on Windows. Mitigated by staggering worker startup delays by 2.0 seconds.",
            "3. Memory Thrashing: Parallel workers saturated RAM. Resolved by limiting PyTorch threads inside workers to 2.",
            "4. Adversarial Candidates: Fake profiles scoring high on semantic search. Mitigated by creating a 4-rule Feature Engineer checker."
        ]
    },
    {
        "title": "Results & Accomplishments",
        "subtitle": "Robust, Validated, Hackathon-Ready Platform",
        "points": [
            "Validated Submissions: Outputs submission.csv containing candidate_id, rank (integer), score (non-increasing), and reasoning.",
            "Format Compliance: Passed all validation checks of validate_submission.py (ties, order, ranks, shapes).",
            "Robust Codebase: Passed 362/362 unit and integration tests with coverage tracking.",
            "Containerized deployment: Fully runnable in multi-container environment via docker-compose."
        ]
    },
    {
        "title": "Future Improvements",
        "subtitle": "Planned Roadmap for Production Upgrades",
        "points": [
            "GPU Acceleration: Integrate automatic CUDA acceleration inside workers for sub-minute precomputations.",
            "Recruiter-in-the-Loop: Leverage active learning to tune scorer weights based on click/hire actions.",
            "Advanced Parsing: Integrate LLMs (e.g. Gemini API) to parse unstructured resumes into structured JSON profiles.",
            "Multi-Lingual Support: Upgrade bi-encoder to multi-lingual models to support global recruitment campaigns."
        ]
    },
    {
        "title": "Streamlit Demo & UI Dashboard",
        "subtitle": "Interactive Dashboard Providing Complete Platform Visibility",
        "points": [
            "Main Panel: Detailed table of Top-100 ranked candidates with search, filter, and CSV export capabilities.",
            "Score Breakdown: Transparent Plotly radar and bar charts comparing bi-encoder, cross-encoder, and heuristic scores.",
            "Skill Alignment: Interactive side-by-side comparison of candidate skills vs JD requirements.",
            "API Docs: Integrated OpenAPI/Swagger UI detailing the /rank, /health, and /top100 endpoints."
        ]
    },
    {
        "title": "Conclusion",
        "subtitle": "Production-Grade AI Talent Intelligence Platform",
        "points": [
            "AI-Powered: Integrates bi-encoders, vector stores, and cross-encoders for maximum search recall and ranking precision.",
            "Transparent & Fair: Provides detailed, human-readable reasonings and excludes adversarial profiles.",
            "High Performance: Optimized to run seamlessly on standard CPU-only hardware under memory and time constraints.",
            "Submission Ready: Repository pushed, validated, and documented for hackathon delivery."
        ]
    }
]


def create_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_slide_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(SLIDE_DATA):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Determine background color: Title & Conclusion are Dark Navy, rest are Light Gray
        background = slide.background
        fill = background.fill
        fill.solid()
        if idx in [0, len(SLIDE_DATA) - 1]:
            fill.fore_color.rgb = DARK_NAVY
            title_color = TEXT_WHITE
            subtitle_color = ACCENT_CYAN
            text_color = TEXT_WHITE
        else:
            fill.fore_color.rgb = LIGHT_GRAY
            title_color = PRIMARY_BLUE
            subtitle_color = SECONDARY_BLUE
            text_color = MUTED_DARK
            
        # Draw Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info["title"]
        p.font.name = "Arial"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        # Draw Subtitle
        if slide_info["subtitle"]:
            p2 = tf.add_paragraph()
            p2.text = slide_info["subtitle"]
            p2.font.name = "Arial"
            p2.font.size = Pt(18)
            p2.font.color.rgb = subtitle_color
            p2.space_before = Pt(5)

        # Draw Points
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5.0))
        c_tf = content_box.text_frame
        c_tf.word_wrap = True
        
        for p_idx, pt in enumerate(slide_info["points"]):
            if p_idx == 0:
                p_obj = c_tf.paragraphs[0]
            else:
                p_obj = c_tf.add_paragraph()
            p_obj.text = pt
            p_obj.font.name = "Arial"
            # Title slide has larger text
            p_obj.font.size = Pt(22) if idx in [0, len(SLIDE_DATA) - 1] else Pt(16)
            p_obj.font.color.rgb = text_color
            p_obj.space_after = Pt(12)
            
            # Simple indentation for sub-bullets
            if pt.startswith("  -") or pt.startswith("    -"):
                p_obj.level = 1
                
    out_path = "AI_Talent_Intelligence_Platform.pptx"
    prs.save(out_path)
    print(f"PowerPoint Presentation saved to: {os.path.abspath(out_path)}")


def create_pdf() -> None:
    pdf_filename = "AI_Talent_Intelligence_Platform.pdf"
    
    # Simple landscape slide-like PDF document
    doc = SimpleDocTemplate(
        pdf_filename,
        pagesize=landscape(letter),
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    style_slide_title_dark = ParagraphStyle(
        'SlideTitleDark',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=32,
        leading=38,
        textColor=colors.HexColor("#FFFFFF"),
        spaceAfter=6
    )
    style_slide_subtitle_dark = ParagraphStyle(
        'SlideSubtitleDark',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=16,
        leading=20,
        textColor=PDF_ACCENT_CYAN,
        spaceAfter=15
    )
    style_slide_text_dark = ParagraphStyle(
        'SlideTextDark',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=15,
        leading=22,
        textColor=colors.HexColor("#E2E8F0"),
        spaceAfter=10
    )
    
    style_slide_title_light = ParagraphStyle(
        'SlideTitleLight',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=32,
        leading=38,
        textColor=PDF_PRIMARY_BLUE,
        spaceAfter=6
    )
    style_slide_subtitle_light = ParagraphStyle(
        'SlideSubtitleLight',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=16,
        leading=20,
        textColor=PDF_SECONDARY_BLUE,
        spaceAfter=15
    )
    style_slide_text_light = ParagraphStyle(
        'SlideTextLight',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=14,
        leading=20,
        textColor=PDF_MUTED_DARK,
        spaceAfter=8
    )

    story = []
    
    for idx, slide_info in enumerate(SLIDE_DATA):
        is_dark = idx in [0, len(SLIDE_DATA) - 1]
        
        # We wrap each slide's content in a full-page Table with specific background
        slide_title_style = style_slide_title_dark if is_dark else style_slide_title_light
        slide_subtitle_style = style_slide_subtitle_dark if is_dark else style_slide_subtitle_light
        slide_text_style = style_slide_text_dark if is_dark else style_slide_text_light
        
        slide_content = []
        slide_content.append(Paragraph(slide_info["title"], slide_title_style))
        if slide_info["subtitle"]:
            slide_content.append(Paragraph(slide_info["subtitle"], slide_subtitle_style))
        slide_content.append(Spacer(1, 15))
        
        for pt in slide_info["points"]:
            # Format sub-bullets nicely in PDF
            indent = ""
            text = pt
            if pt.startswith("  -"):
                indent = "&nbsp;&nbsp;&nbsp;&nbsp;&bull;&nbsp;"
                text = pt[3:]
            elif pt.startswith("    -"):
                indent = "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&bull;&nbsp;"
                text = pt[5:]
            else:
                indent = "&bull;&nbsp;"
            
            p_text = f"{indent}{text}"
            slide_content.append(Paragraph(p_text, slide_text_style))
            
        # Wrap everything in a single table cell to fill background
        slide_table = Table([[slide_content]], colWidths=[700], rowHeights=[500])
        
        bg_color = PDF_DARK_NAVY if is_dark else PDF_LIGHT_GRAY
        
        slide_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), bg_color),
            ('TOPPADDING', (0,0), (-1,-1), 30),
            ('BOTTOMPADDING', (0,0), (-1,-1), 30),
            ('LEFTPADDING', (0,0), (-1,-1), 40),
            ('RIGHTPADDING', (0,0), (-1,-1), 40),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ]))
        
        story.append(slide_table)
        if idx < len(SLIDE_DATA) - 1:
            story.append(PageBreak())
            
    doc.build(story)
    print(f"PDF Presentation saved to: {os.path.abspath(pdf_filename)}")


if __name__ == "__main__":
    create_pptx()
    create_pdf()
